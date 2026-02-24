"""PPTX to Markdown converter using python-pptx."""

from pathlib import Path

from .base import BaseConverter, ConversionResult
from utils.markdown_clean import clean_markdown


class PptxConverter(BaseConverter):
    supported_extensions = (".pptx",)

    def convert(self, source: str | Path, **kwargs) -> ConversionResult:
        path = Path(source)

        from pptx import Presentation

        prs = Presentation(str(path))
        md_parts = [f"# {path.name}\n"]

        for i, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide)
            heading = f"## Slide {i}" + (f": {title}" if title else "")
            md_parts.append(heading)

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            md_parts.append(f"- {text}")

                if shape.has_table:
                    md_parts.append(self._table_to_markdown(shape.table))

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    md_parts.append(f"\n> **Notes:** {notes}")

            md_parts.append("")  # blank line between slides

        content = clean_markdown("\n".join(md_parts))
        return ConversionResult(content=content, title=path.name, source=str(path))

    def _get_slide_title(self, slide) -> str | None:
        """Extract the title text from a slide."""
        if slide.shapes.title:
            return slide.shapes.title.text.strip() or None
        return None

    def _table_to_markdown(self, table) -> str:
        """Convert a PPTX table to a markdown table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if len(rows) >= 1:
            # Add separator after header
            col_count = len(table.rows[0].cells)
            separator = "| " + " | ".join("---" for _ in range(col_count)) + " |"
            rows.insert(1, separator)

        return "\n".join(rows)
